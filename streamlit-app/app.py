import streamlit as st
import pandas as pd
import PyPDF2
import docx
import numpy as np
from PIL import Image
import pytesseract
import google.generativeai as genai
from pptx import Presentation
import re
import json
from typing import List, Dict, Any

# ---------------- Config ----------------
st.set_page_config(page_title="Smart Document Assistant", layout="wide")

# Initialize Gemini
if "gemini_client" not in st.session_state:
    api_key = st.text_input("Enter your Gemini API Key:", type="password", key="api_key_input")
    if api_key:
        try:
            genai.configure(api_key=api_key)
            st.session_state.gemini_client = genai.GenerativeModel('gemini-2.0-flash-exp')
            st.success("✅ API Key configured successfully!")
        except Exception as e:
            st.error(f"❌ API Key error: {str(e)}")
    else:
        st.info("👆 Please enter your Gemini API key to continue")
        st.stop()

# Path to Tesseract (adjust if needed)
try:
    pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except:
    pass

# ---------------- Session State ----------------
if 'documents' not in st.session_state:
    st.session_state.documents = {}
if 'conversation_history' not in st.session_state:
    st.session_state.conversation_history = []
if 'current_files' not in st.session_state:
    st.session_state.current_files = []

# ---------------- Helper Functions ----------------

def clean_text(text: str) -> str:
    """Clean and normalize text"""
    if not text:
        return ""
    # Remove extra whitespace, keep single spaces
    text = re.sub(r'\s+', ' ', text.strip())
    return text

def extract_text_from_file(file) -> Dict[str, Any]:
    """Extract text from various file formats"""
    filename = file.name
    file_extension = filename.lower().split('.')[-1]
    
    result = {
        'filename': filename,
        'type': file_extension.upper(),
        'content': "",
        'structured_data': None,
        'raw_content': ""
    }
    
    try:
        if file_extension == 'pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            pages = []
            for i, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text.strip():
                    pages.append(f"[Page {i+1}]\n{clean_text(page_text)}")
            result['content'] = "\n\n".join(pages)
            result['raw_content'] = result['content']
            
        elif file_extension == 'docx':
            doc = docx.Document(file)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(clean_text(para.text))
            result['content'] = "\n\n".join(paragraphs)
            result['raw_content'] = result['content']
            
        elif file_extension == 'txt':
            content = file.read().decode('utf-8')
            result['content'] = clean_text(content)
            result['raw_content'] = result['content']
            
        elif file_extension == 'csv':
            df = pd.read_csv(file)
            # Store original dataframe
            result['structured_data'] = df
            
            # Create detailed text representation
            content_parts = []
            content_parts.append(f"CSV File: {filename}")
            content_parts.append(f"Total Rows: {len(df)}")
            content_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
            content_parts.append("\nDetailed Data:")
            
            # Add each row with clear formatting
            for idx, row in df.iterrows():
                row_text = f"Row {idx + 1}:"
                for col in df.columns:
                    value = str(row[col]).strip() if pd.notna(row[col]) else ""
                    if value:
                        row_text += f" {col}='{value}',"
                content_parts.append(row_text.rstrip(','))
            
            result['content'] = "\n".join(content_parts)
            result['raw_content'] = df.to_string(index=True)
            
        elif file_extension in ['xlsx', 'xls']:
            excel_file = pd.ExcelFile(file)
            all_sheets = {}
            content_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file, sheet_name=sheet_name)
                all_sheets[sheet_name] = df
                
                content_parts.append(f"Sheet: {sheet_name}")
                content_parts.append(f"Rows: {len(df)}, Columns: {', '.join(df.columns.tolist())}")
                
                # Add each row with clear formatting
                for idx, row in df.iterrows():
                    row_text = f"Row {idx + 1}:"
                    for col in df.columns:
                        value = str(row[col]).strip() if pd.notna(row[col]) else ""
                        if value:
                            row_text += f" {col}='{value}',"
                    content_parts.append(row_text.rstrip(','))
                content_parts.append("")
            
            result['structured_data'] = all_sheets
            result['content'] = "\n".join(content_parts)
            result['raw_content'] = "\n\n".join([f"Sheet: {name}\n{df.to_string(index=True)}" for name, df in all_sheets.items()])
            
        elif file_extension in ['pptx', 'ppt']:
            prs = Presentation(file)
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(clean_text(shape.text))
                if slide_content:
                    slides.append(f"[Slide {i+1}]\n" + "\n".join(slide_content))
            result['content'] = "\n\n".join(slides)
            result['raw_content'] = result['content']
            
        elif file_extension in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            image = Image.open(file)
            
            # Try OCR
            try:
                ocr_text = pytesseract.image_to_string(image)
                ocr_text = clean_text(ocr_text)
            except Exception as ocr_error:
                ocr_text = f"OCR Error: {str(ocr_error)}"
            
            # Image description
            image_info = f"""Image File: {filename}
Format: {image.format}
Size: {image.size[0]}x{image.size[1]} pixels
Mode: {image.mode}

Extracted Text (OCR):
{ocr_text if ocr_text and len(ocr_text.strip()) > 3 else "No clear text found in image"}

Image Content: This is an image file that may contain visual information not captured by OCR."""
            
            result['content'] = image_info
            result['raw_content'] = ocr_text
            result['structured_data'] = image
            
    except Exception as e:
        result['content'] = f"Error processing {filename}: {str(e)}"
        result['raw_content'] = result['content']
    
    return result

def search_in_documents(query: str, documents: Dict) -> str:
    """Search for relevant content in documents with improved matching"""
    query_lower = query.lower()
    query_words = set(re.findall(r'\b\w+\b', query_lower))
    
    relevant_content = []
    
    # Debug: Show what we're searching for
    print(f"DEBUG: Query words: {query_words}")
    
    for filename, doc_data in documents.items():
        content = doc_data['content'].lower()
        
        # Debug: Show document content length
        print(f"DEBUG: Document {filename} content length: {len(doc_data['content'])}")
        print(f"DEBUG: First 200 chars: {doc_data['content'][:200]}")
        
        # Check if any query words appear in content
        content_words = set(re.findall(r'\b\w+\b', content))
        word_matches = len(query_words.intersection(content_words))
        
        print(f"DEBUG: Word matches for {filename}: {word_matches}")
        
        # More lenient matching - if document has content, include it for broad questions
        if word_matches > 0 or (len(query_words) <= 3 and len(doc_data['content']) > 50):
            # Add the full content for this document
            relevant_content.append(f"=== From {filename} ===\n{doc_data['content']}")
            
            # Also add raw content if different (for structured data)
            if doc_data['raw_content'] != doc_data['content']:
                relevant_content.append(f"=== Raw Data from {filename} ===\n{doc_data['raw_content']}")
    
    # If no matches found but we have documents, return all content for very general questions
    if not relevant_content and documents:
        general_questions = ['what', 'about', 'document', 'content', 'summary', 'summarize', 'tell', 'describe']
        if any(word in query_lower for word in general_questions):
            print("DEBUG: Using fallback - returning all document content")
            for filename, doc_data in documents.items():
                relevant_content.append(f"=== From {filename} ===\n{doc_data['content']}")
    
    result = "\n\n".join(relevant_content) if relevant_content else ""
    print(f"DEBUG: Final context length: {len(result)}")
    return result

def calculate_from_structured_data(query: str, documents: Dict) -> str:
    """Perform calculations on structured data (Excel/CSV)"""
    query_lower = query.lower()
    calculations = []
    
    for filename, doc_data in documents.items():
        if doc_data['structured_data'] is not None:
            if isinstance(doc_data['structured_data'], pd.DataFrame):
                df = doc_data['structured_data']
                calculations.append(f"=== Calculations from {filename} ===")
                calculations.append(f"Total rows: {len(df)}")
                
                # Try to find numerical columns for calculations
                for col in df.columns:
                    if df[col].dtype in ['int64', 'float64']:
                        total = df[col].sum()
                        mean_val = df[col].mean()
                        calculations.append(f"{col}: Sum={total}, Average={mean_val:.2f}")
                
            elif isinstance(doc_data['structured_data'], dict):  # Excel with multiple sheets
                for sheet_name, df in doc_data['structured_data'].items():
                    calculations.append(f"=== Calculations from {filename} - Sheet: {sheet_name} ===")
                    calculations.append(f"Total rows: {len(df)}")
                    
                    for col in df.columns:
                        if df[col].dtype in ['int64', 'float64']:
                            total = df[col].sum()
                            mean_val = df[col].mean()
                            calculations.append(f"{col}: Sum={total}, Average={mean_val:.2f}")
    
    return "\n".join(calculations) if calculations else ""

def generate_answer(query: str, context: str, documents: Dict) -> str:
    """Generate answer using Gemini with strict guardrails"""
    
    if not context.strip():
        return "I don't know - this information is not available in the uploaded documents."
    
    # Enhanced context with calculations for structured data
    enhanced_context = context
    if any('total' in query.lower() or 'sum' in query.lower() or 'calculate' in query.lower() or 'count' in query.lower() for _ in [1]):
        calc_info = calculate_from_structured_data(query, documents)
        if calc_info:
            enhanced_context += f"\n\n=== CALCULATIONS ===\n{calc_info}"
    
    # Strict prompt with multiple guardrails
    prompt = f"""You are a document analysis assistant. You must ONLY use information from the provided documents to answer questions.

STRICT RULES - FOLLOW EXACTLY:
1. ONLY use information explicitly stated in the context below
2. Do NOT use any external knowledge or general information
3. If the answer is not clearly in the context, respond EXACTLY with: "I don't know - this information is not available in the uploaded documents."
4. For numerical questions, be precise and check ALL data entries carefully
5. For comparisons between files, look at ALL provided documents
6. For image descriptions, ONLY describe what you can see from the OCR text or image information provided
7. Be case-insensitive when searching for information
8. When counting or calculating, examine EVERY single entry in the data
9. If asked about something not in the documents, always say "I don't know"

CONTEXT FROM UPLOADED DOCUMENTS:
{enhanced_context}

QUESTION: {query}

Remember: If you cannot find the answer in the context above, you MUST respond with "I don't know - this information is not available in the uploaded documents."

ANSWER:"""

    try:
        response = st.session_state.gemini_client.generate_content(
            prompt,
            generation_config=genai.types.GenerationConfig(
                temperature=0.1,
                top_p=0.8,
                top_k=20,
                max_output_tokens=1000,
            ),
        )
        
        answer = response.text.strip()
        
        # Additional validation
        if not answer or len(answer) < 5:
            return "I don't know - this information is not available in the uploaded documents."
        
        # Check if the answer seems to use external knowledge
        external_indicators = [
            "generally", "typically", "usually", "in general", "commonly",
            "it is known", "studies show", "research indicates", "experts say",
            "according to", "based on my knowledge"
        ]
        
        answer_lower = answer.lower()
        for indicator in external_indicators:
            if indicator in answer_lower:
                return "I don't know - this information is not available in the uploaded documents."
        
        return answer
        
    except Exception as e:
        return f"Error generating response: {str(e)}"

# ---------------- Streamlit UI ----------------

st.title("🤖 Smart Document Assistant")
st.markdown("Upload documents and ask questions - I'll only answer based on your uploaded content!")

# Sidebar for file management
with st.sidebar:
    st.header("📁 Document Management")
    
    uploaded_files = st.file_uploader(
        "Upload Documents",
        type=["pdf", "docx", "txt", "csv", "xlsx", "xls", "pptx", "ppt", "png", "jpg", "jpeg", "tiff", "bmp"],
        accept_multiple_files=True,
        key="file_uploader"
    )
    
    if uploaded_files:
        if st.button("🔄 Process Documents", type="primary"):
            with st.spinner("Processing documents..."):
                st.session_state.documents = {}
                st.session_state.current_files = []
                
                for file in uploaded_files:
                    doc_data = extract_text_from_file(file)
                    st.session_state.documents[file.name] = doc_data
                    st.session_state.current_files.append(file.name)
                
                # Reset conversation when new documents are uploaded
                st.session_state.conversation_history = []
                st.success(f"✅ Processed {len(uploaded_files)} documents!")
    
    # Show processed documents
    if st.session_state.current_files:
        st.subheader("📋 Loaded Documents")
        for filename in st.session_state.current_files:
            doc = st.session_state.documents[filename]
            with st.expander(f"{doc['type']}: {filename}"):
                st.write(f"Content length: {len(doc['content'])} chars")
                if doc['structured_data'] is not None:
                    if isinstance(doc['structured_data'], pd.DataFrame):
                        st.write(f"Data shape: {doc['structured_data'].shape}")
                        st.dataframe(doc['structured_data'].head())
                    elif isinstance(doc['structured_data'], dict):
                        st.write(f"Sheets: {len(doc['structured_data'])}")

# Main chat interface
if not st.session_state.documents:
    st.info("👆 Please upload some documents to get started!")
else:
    # Chat container
    chat_container = st.container()
    
    with chat_container:
        # Display conversation history
        for i, (question, answer) in enumerate(st.session_state.conversation_history):
            with st.chat_message("user"):
                st.write(f"**Q{i+1}:** {question}")
            with st.chat_message("assistant"):
                st.write(f"**A{i+1}:** {answer}")
    
    # Input area
    col1, col2, col3 = st.columns([6, 1, 1])
    
    with col1:
        user_question = st.text_input(
            "Ask me anything about your documents:",
            placeholder="e.g., What is the total quantity of product X? Compare these files...",
            key="user_input"
        )
    
    with col2:
        ask_button = st.button("🚀 Ask", type="primary", use_container_width=True)
    
    with col3:
        clear_button = st.button("🗑️ Clear", use_container_width=True)
    
    # Handle user input
    if ask_button and user_question.strip():
        with st.spinner("🔍 Searching through your documents..."):
            # Search for relevant content
            context = search_in_documents(user_question, st.session_state.documents)
            
            # Debug info
            with st.expander("🔧 Debug Info (click to expand)"):
                st.write(f"**Query:** {user_question}")
                st.write(f"**Context found:** {len(context)} characters")
                if context:
                    st.write("**Context preview:**")
                    st.text(context[:500] + "..." if len(context) > 500 else context)
                else:
                    st.write("**No context found!**")
                    st.write("**Available documents:**")
                    for filename, doc_data in st.session_state.documents.items():
                        st.write(f"- {filename}: {len(doc_data['content'])} chars")
                        st.write(f"  Preview: {doc_data['content'][:100]}...")
            
            # Generate answer
            answer = generate_answer(user_question, context, st.session_state.documents)
            
            # Add to conversation history
            st.session_state.conversation_history.append((user_question, answer))
            
            # Display the new Q&A
            with st.chat_message("user"):
                st.write(f"**Q{len(st.session_state.conversation_history)}:** {user_question}")
            with st.chat_message("assistant"):
                st.write(f"**A{len(st.session_state.conversation_history)}:** {answer}")
            
            # Clear the input
            st.rerun()
    
    # Clear conversation
    if clear_button:
        st.session_state.conversation_history = []
        st.rerun()
    
    # Download conversation
    if st.session_state.conversation_history:
        conversation_text = ""
        for i, (q, a) in enumerate(st.session_state.conversation_history):
            conversation_text += f"Q{i+1}: {q}\nA{i+1}: {a}\n\n"
        
        st.download_button(
            "📥 Download Conversation",
            conversation_text,
            "chat_history.txt",
            "text/plain",
            use_container_width=True
        )

# Footer
st.markdown("---")
st.markdown("💡 **Tips:** Ask specific questions, compare documents, request calculations from Excel/CSV files, or ask about image content!")

# Add some example questions based on uploaded files
if st.session_state.documents:
    with st.expander("💭 Example Questions"):
        st.markdown("**Click to copy these example questions:**")
        examples = [
            "What is this document about?",
            "Summarize the main points from all documents",
            "What is the total quantity in the data?",
            "Compare the information between the uploaded files",
            "What does the image contain?",
            "Show me all entries for [specific item]",
            "Calculate the sum of [column name]",
            "Who appears most frequently in the data?"
        ]
        for example in examples:
            st.code(example)
